import sys
import subprocess
import os
from PIL import Image

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing via pip...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def main():
    prs = Presentation()
    
    # 16:9 Widescreen slide setup
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # WEB APP DESIGN SYSTEM (AICA Slate Dark Mode Theme)
    COLOR_BG = RGBColor(15, 23, 42)          # #0f172a (Deep Dark Slate backdrop)
    COLOR_CARD = RGBColor(30, 41, 59)        # #1e293b (Slate Card Background)
    COLOR_CARD_ALT = RGBColor(15, 23, 42)    # #0f172a (Inner Dark contrast)
    COLOR_BORDER = RGBColor(51, 65, 85)      # #334155 (Slate Card Border)
    COLOR_TEXT_MAIN = RGBColor(255, 255, 255)# Pure White (#ffffff)
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Gray Muted (#94a3b8)
    COLOR_BLUE = RGBColor(56, 189, 248)      # #38bdf8 (Vibrant Cyan/Blue Accent)
    COLOR_INDIGO = RGBColor(99, 102, 241)    # #6366f1 (Indigo Accent)
    COLOR_GREEN = RGBColor(16, 185, 129)     # #10b981 (Success Mint Green)
    COLOR_RED = RGBColor(244, 63, 94)        # #f43f5e (Warning Red/Rose)

    FONT_NAME = "Noto Sans KR"
    blank_layout = prs.slide_layouts[6]
    img_dir = r"C:\Users\ATECCN\.gemini\antigravity\brain\251e0be8-b054-44a9-82c7-d28a8df26ff8"

    # Helper: Set dark background color
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    # Helper: Draw header bar (AICA Dark Mode style)
    def add_header(slide, title, category="공통", slide_num_str=""):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.4))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_BLUE
        p_cat.font.name = FONT_NAME
        
        # Main Title (White bold style)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(11), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = FONT_NAME
        
        # Header separator line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER
        line.line.fill.background()
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = f"AICA 자격심사 시스템 사용자 매뉴얼 · Web App Dark Theme"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = COLOR_TEXT_MUTED
        p_foot.font.name = FONT_NAME
        
        if slide_num_str:
            num_box = slide.shapes.add_textbox(Inches(11.833), Inches(6.9), Inches(0.9), Inches(0.4))
            tf_num = num_box.text_frame
            p_num = tf_num.paragraphs[0]
            p_num.alignment = PP_ALIGN.RIGHT
            p_num.text = slide_num_str
            p_num.font.size = Pt(9)
            p_num.font.color.rgb = COLOR_TEXT_MUTED
            p_num.font.name = FONT_NAME

    # Helper: Create text card (Sleek Slate card style)
    def add_card_box(slide, left, top, width, height, title, bullets, is_alt=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD_ALT if is_alt else COLOR_CARD
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(1)
        
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BLUE if is_alt else COLOR_TEXT_MAIN
        p_title.font.name = FONT_NAME
        p_title.space_after = Pt(12)
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_MUTED
            p.font.name = FONT_NAME
            p.space_after = Pt(6)

    # Helper: Insert screenshot image with Slate frame
    def add_screenshot_frame(slide, left, top, width, height, image_name):
        img_path = os.path.join(img_dir, image_name)
        if not os.path.exists(img_path):
            add_card_box(slide, left, top, width, height, "화면 스크린샷 준비 중", [
                "해당 기능 설명에 부합하는 앱 화면 캡처 파일이 누락되었습니다.",
                f"파일명: {image_name}",
                "이후 캡처본을 배치하면 자동 갱신됩니다."
            ], is_alt=True)
            return

        # Frame container (Slate card rounded box style with glowing indigo border)
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = COLOR_CARD
        frame.line.color.rgb = COLOR_INDIGO
        frame.line.width = Pt(1.5)
        
        # Center image within the frame keeping aspect ratio
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
            
            aspect = img_w / img_h
            box_w = width - Inches(0.3)
            box_h = height - Inches(0.3)
            box_aspect = box_w / box_h
            
            if aspect > box_aspect:
                w = box_w
                h = box_w / aspect
                x = left + Inches(0.15)
                y = top + Inches(0.15) + (box_h - h) / 2
            else:
                h = box_h
                w = box_h * aspect
                x = left + Inches(0.15) + (box_w - w) / 2
                y = top + Inches(0.15)
            
            slide.shapes.add_picture(img_path, x, y, width=w, height=h)
        except Exception as e:
            print(f"Error rendering image {image_name}: {e}")

    # Helper: Add standard two-column slide layout
    def add_two_column_slide(slide, title, left_card_title, left_bullets, right_image_name, category="공통", slide_num=""):
        add_header(slide, title, category, slide_num)
        
        # Left column: Text card
        add_card_box(slide, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.5), left_card_title, left_bullets, is_alt=False)
        
        # Right column: Screenshot
        add_screenshot_frame(slide, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.5), right_image_name)

    # =========================================================================
    # SLIDE 1: Title Slide (Cover Page - Web App style)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)
    
    # Left Title block
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p_meta = tf1.paragraphs[0]
    p_meta.text = "AICA · AI Certification for ATEC"
    p_meta.font.size = Pt(13)
    p_meta.font.bold = True
    p_meta.font.color.rgb = COLOR_BLUE
    p_meta.font.name = FONT_NAME
    p_meta.space_after = Pt(20)
    
    p_title = tf1.add_paragraph()
    p_title.text = "자격심사 시스템\n사용자 매뉴얼"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    p_title.font.name = FONT_NAME
    p_title.space_after = Pt(20)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "운영간사 · 심사위원용"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.font.name = FONT_NAME
    
    # Left Footer Meta
    tb_meta_foot = slide1.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(6.8), Inches(0.8))
    tf_meta_foot = tb_meta_foot.text_frame
    p_mf = tf_meta_foot.paragraphs[0]
    p_mf.text = "2026. 07\nhttps://aica-competency.web.app/"
    p_mf.font.size = Pt(11)
    p_mf.font.color.rgb = COLOR_TEXT_MUTED
    p_mf.font.name = FONT_NAME
    
    # Right Dark Card graphic
    rcard = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.0), Inches(4.5), Inches(5.5))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = COLOR_CARD
    rcard.line.color.rgb = COLOR_BORDER
    rcard.line.width = Pt(1.5)
    
    # Rcard contents
    tb_rcard = slide1.shapes.add_textbox(Inches(8.2), Inches(1.5), Inches(4.1), Inches(4.5))
    tf_rcard = tb_rcard.text_frame
    tf_rcard.word_wrap = True
    
    p_rc_badge = tf_rcard.paragraphs[0]
    p_rc_badge.alignment = PP_ALIGN.CENTER
    p_rc_badge.text = "AICA"
    p_rc_badge.font.size = Pt(28)
    p_rc_badge.font.bold = True
    p_rc_badge.font.color.rgb = COLOR_TEXT_MAIN
    p_rc_badge.font.name = FONT_NAME
    p_rc_badge.space_after = Pt(30)
    
    p_rc_title = tf_rcard.add_paragraph()
    p_rc_title.alignment = PP_ALIGN.CENTER
    p_rc_title.text = "공정한 3인 패널 심사"
    p_rc_title.font.size = Pt(18)
    p_rc_title.font.bold = True
    p_rc_title.font.color.rgb = COLOR_TEXT_MAIN
    p_rc_title.font.name = FONT_NAME
    p_rc_title.space_after = Pt(10)
    
    p_rc_sub = tf_rcard.add_paragraph()
    p_rc_sub.alignment = PP_ALIGN.CENTER
    p_rc_sub.text = "현업 · 기술 · 보안"
    p_rc_sub.font.size = Pt(14)
    p_rc_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_rc_sub.font.name = FONT_NAME
    p_rc_sub.space_after = Pt(40)
    
    p_rc_flow = tf_rcard.add_paragraph()
    p_rc_flow.alignment = PP_ALIGN.CENTER
    p_rc_flow.text = "등록 ➔ 검토 ➔ 채점 ➔ 판정"
    p_rc_flow.font.size = Pt(14)
    p_rc_flow.font.bold = True
    p_rc_flow.font.color.rgb = COLOR_BLUE
    p_rc_flow.font.name = FONT_NAME

    # =========================================================================
    # SLIDE 2: Process Overview (공통 - 두 역할이 이어져야 한 건의 심사가 완료됩니다)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_two_column_slide(
        slide2,
        "두 역할이 이어져야 한 건의 심사가 완료됩니다",
        "자격 심사 전체 워크플로우",
        [
            "운영간사와 심사위원 간의 원활한 업무 연계 흐름:",
            "• 1단계 (과제 등록): 간사가 지원자명, 과제설명, 증빙 자료 등록",
            "• 2단계 (3인 패널 배정): 현업, 기술, 보안별 심사위원 자동 매핑",
            "• 3단계 (독립 심사): 각 위원이 증빙을 검증하고 5단계 평가 제출",
            "• 4단계 (종합 판정): 3인 점수 평균 및 결격 검증 후 판정서 발급",
            "\n💡 우측 화면: 리액트 및 클라우드 파이어스토어(Firestore) 데이터베이스 기반 실시간 평가 시스템의 전체 구조"
        ],
        "aica_evaluation_workspace_mockup_1784167999282.png",
        "공통",
        "02"
    )

    # =========================================================================
    # SLIDE 3: Login (공통 - 성명과 사내 이메일로 로그인합니다)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_two_column_slide(
        slide3,
        "성명과 사내 이메일로 로그인합니다",
        "로그인 순서 및 계정 인증",
        [
            "1. 성명(ID)에 등록된 성명을 입력",
            "2. 이메일(Password)에 사내 이메일 입력",
            "3. 「시스템 로그인」 선택",
            "4. 상단 역할 표시가 본인 권한인지 확인",
            "\n* 입력이 없거나 일치하지 않으면 로그인 오류가 표시됩니다. 운영간사에게 등록 정보 확인을 요청하세요.",
            "\n오른쪽 화면: 실제 리액트 시스템 로그인 창 및 하단 테스트 계정 정보창 캡처 화면"
        ],
        "media__1784176101924.png",
        "공통",
        "03"
    )

    # =========================================================================
    # SLIDE 4: Roles Description (공통 - 먼저 내 역할과 완료 조건을 확인하세요)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_two_column_slide(
        slide4,
        "먼저 내 역할과 완료 조건을 확인하세요",
        "사용자 역할 분리 및 완료 기준",
        [
            "운영간사의 핵심 임무 및 완료 조건:",
            "• 신규 평가과제 및 후보자 등록",
            "• 3인 패널 지정 및 이해상충 자가검증",
            "• 큐·상태·점수 진척도 실시간 모니터링",
            "• 완료 신호: 3인의 심사가 완료되고 최종 판정서 생성",
            "\n심사위원의 핵심 임무 및 완료 조건:",
            "• 제출물·시연·링크 검토 및 공통 보안 체크리스트 확인",
            "• Level 3 / Level 4 루브릭 평가 및 면접 질문 작성",
            "• 완료 신호: 내 심사 상태 배지가 「완료」로 표시됨",
            "\n오른쪽 화면: 실시간 진행 상황을 요약하는 관리용 대시보드 인터페이스"
        ],
        "media__1784188798650.png",
        "공통",
        "04"
    )

    # =========================================================================
    # SLIDE 5: Admin Overview (운영간사 - 대시보드에서 전체 심사 현황부터 확인합니다)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_two_column_slide(
        slide5,
        "대시보드에서 전체 심사 현황부터 확인합니다",
        "AICA 자격심사 관리 & 대기 큐 구조",
        [
            "1. 시뮬레이션 데이터 불러오기: 가상의 테스트 데이터셋(후보자 9명 및 3인 패널 평가 결과) 복원",
            "2. 데이터 초기화: 실무 운영 착수를 위해 전체 평가 과제 및 결과를 깨끗하게 초기화",
            "3. 요약 지표: 전체, 완료, 진행 중, 대기 인원을 실시간 확인 (클릭 시 명단 모달 팝업)",
            "4. 관계사 평균: 완료된 심사를 기준으로 계열사별 평균 점수를 모니터링",
            "5. 전체 큐: 필터와 액션을 통한 개별 심사 통합 관리",
            "\n💡 대시보드 우측 상단 배지를 통해 테스트 시연용 가상 데이터 복원 및 실무용 초기화를 손쉽게 제어할 수 있습니다.",
            "\n오른쪽 화면: 운영간사용 메인 관리 대시보드 및 평가 큐 상세 화면"
        ],
        "media__1784188798650.png",
        "운영간사",
        "05"
    )

    # =========================================================================
    # SLIDE 6: Register Candidate (운영간사 - 신규 과제는 필수 정보와 증빙을 한 번에 등록합니다)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_two_column_slide(
        slide6,
        "신규 과제는 필수 정보와 증빙을 한 번에 등록합니다",
        "후보자 정보 및 증빙 링크 등록 폼",
        [
            "01 지원자 정보:",
            "• 성명, 사내 이메일, 소속 관계사/부서, 자격 레벨(L3/L4), 과제 유형",
            "\n02 과제 내용 및 6S 요약:",
            "• 평가 대상 과제명, Pain Point, AI 해결 방안, 핵심 기대효과 기술",
            "\n03 증빙 파일 링크 등록:",
            "• 디렉토리 구조 요약, 결과 보고서 PDF 주소, Git 소스코드 저장소, 시연 동영상 URL",
            "\n오른쪽 화면: 수동 과제 등록 및 패널 배정 기능이 활성화되는 대시보드 화면"
        ],
        "media__1784188798650.png",
        "운영간사",
        "06"
    )

    # =========================================================================
    # SLIDE 7: Panel Assignment (운영간사 - 패널은 전문분야별 1인씩, 이해상충 없이 배정합니다)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_two_column_slide(
        slide7,
        "패널은 전문분야별 1인씩, 이해상충 없이 배정합니다",
        "심사 위원회 패널 구성 및 최적 추천",
        [
            "1. 현업·사업성 위원 (가치성 검증): 업무 프로세스 개선율 및 사업 효과성 채점",
            "2. AI·기술 위원 (구현성 검증): 기술적 난이도 및 구현 완성도 채점",
            "3. 보안·거버넌스 위원 (신뢰성 검증): 데이터 활용 보안 및 내부 가이드 채점",
            "\n⚡ 최적 패널 자동 추천 로직:",
            "• 피평가자와 동일한 사명 소속 위원을 자동 배제하여 후보군을 우선 순위 정렬합니다.",
            "\n⚠️ 간사가 위원을 직접 수동 선택하는 경우, 소속이 일치하는 위원 옆에는 빨간색 [이해상충 주의] 레이블이 노출됩니다.",
            "\n오른쪽 화면: 대시보드 우측의 위원 배정 큐 영역"
        ],
        "media__1784188798650.png",
        "운영간사",
        "07"
    )

    # =========================================================================
    # SLIDE 8: Queue & Filters (운영간사 - 필터와 상태값으로 병목을 빠르게 찾습니다)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_two_column_slide(
        slide8,
        "필터와 상태값으로 병목을 빠르게 찾습니다",
        "진행 상태 분기 및 필터 모니터링",
        [
            "통합 대기 큐 및 필터 시스템 활용법:",
            "• 필터 종류: 관계사별, 인증 레벨별, 평가 진행 상태별 필터 제공",
            "• 1. 평가대기: 패널 배정 완료 후 아직 위원이 채점을 시작하지 않음",
            "• 2. 심사중: 1명 이상의 위원이 부분 채점을 기록하거나 제출함",
            "• 3. 완료: 3인 패널이 모두 최종 평가를 제출하여 종합 평균이 집계됨",
            "• 4. 보고서 보기: [결과 보기] 액션 버튼을 눌러 종합판정서 조회 및 인쇄",
            "\n💡 팁: '타 위원 진척도' 항목에서 실시간 배지 도트를 확인해 미제출 위원을 즉시 판별할 수 있습니다.",
            "\n오른쪽 화면: 간사용 대기 큐 필터 및 위원 진척 표시 영역"
        ],
        "media__1784188798650.png",
        "운영간사",
        "08"
    )

    # =========================================================================
    # SLIDE 9: Final Report (운영간사 - 3인 제출 후 종합판정서가 완성됩니다)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_two_column_slide(
        slide9,
        "3인 제출 후 종합판정서가 완성됩니다",
        "종합판정 보고서 발급 및 명패 동기화",
        [
            "종합판정서 구성 항목 및 검증 가이드:",
            "• 지원자 기본 정보 및 종합 평균 점수 시각화",
            "• 3인 배정 위원의 세부 루브릭 채점 점수 및 종합 정성 평가 의견 노출",
            "• 하단 서명란: 배정된 실제 심사위원 성명으로 명패 자동 변경 완료",
            "• 출력 최적화: [결과 보고서 인쇄 (PDF)] 버튼을 클릭하여 A4 최적화 보고서 출력 및 사내 공유용 저장 가능",
            "\n⚠️ 결격 사유 발생 시:",
            "• 표절, 중요 기밀 및 API Key 노출 등이 발견되어 결격 처리된 경우, 평균 점수가 70점 이상이어도 최종 판정은 '불합격' 처리됩니다.",
            "\n오른쪽 화면: 서명이 실제 심사위원 성명으로 연동된 결과 보고서 하단부"
        ],
        "media__1784596287593.png",
        "운영간사",
        "09"
    )

    # =========================================================================
    # SLIDE 10: Reviewer Queue (심사위원 - 내 배정 과제 큐에서 심사 우선순위를 잡습니다)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)
    add_two_column_slide(
        slide10,
        "내 배정 과제 큐에서 심사 우선순위를 잡습니다",
        "심사 위원 배정 큐 모니터링",
        [
            "심사위원 전용 배정 과제 리스트 가이드:",
            "• 성명 및 이메일 로그인 시, 본인이 패널로 등록된 평가 과제만 화면에 정렬됩니다.",
            "• 테이블 컬럼: 지원자 정보, 소속 법인/부서, 자격 인증 레벨, 내 평가 진행 현황(대기/완료), 타 위원 제출 진척 상태 배지 표시",
            "\n[심사하기/수정]:",
            "• 채점을 시작하거나 이미 제출한 내용을 재조정할 때 버튼을 선택해 작업장(Workspace)으로 진입합니다.",
            "\n오른쪽 화면: 실제 심사위원이 접하는 과제 목록 테이블 영역"
        ],
        "media__1784176902473.png",
        "심사위원",
        "10"
    )

    # =========================================================================
    # SLIDE 11: Candidate Materials (심사위원 - 채점 전에 제출물의 연결성과 재현 가능성을 봅니다)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11)
    add_two_column_slide(
        slide11,
        "채점 전에 제출물의 연결성과 재현 가능성을 봅니다",
        "3대 증빙 자료 다차원 교차 검증",
        [
            "제출물 상세 탭(Tab) 활용 방식:",
            "• 보고서 요약 (Tab 1): 과제 제안 배경 및 Pain Point 극복 방안, 6S 요약 검토 (문제-해결-기대효과의 논리적 일관성 체크)",
            "• 소스코드 구조 (Tab 2): 디렉토리 트리 구조, 라이브러리 구성 및 API Key 노출 방지 대책 검토",
            "• 시연·링크 (Tab 3): 결과 보고서 PDF 원본, Git 저장소 주소, 구동 동영상 링크 유효성 검증",
            "\n💡 소스코드 및 구동 영상 링크가 손상되었을 경우 정성 의견란에 구체적인 비작동 사유를 기재해 주세요.",
            "\n오른쪽 화면: 실제 심사 작업장 좌측의 3단 증빙 자료 리더(Reader)"
        ],
        "aica_evaluation_workspace_mockup_1784167999282.png",
        "심사위원",
        "11"
    )

    # =========================================================================
    # SLIDE 12: Security Checklist (심사위원 - 공통 보안 체크리스트는 모든 심사의 선행 조건입니다)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12)
    add_two_column_slide(
        slide12,
        "공통 보안 체크리스트는 모든 심사의 선행 조건입니다",
        "5대 공통 보안성 자가진단 규칙",
        [
            "모든 심사의 최우선 선행 체크리스트:",
            "• 1. 개인정보(PII): 이름, 메일주소, 연락처 등 완벽 마스킹 여부",
            "• 2. 기밀 정보 익명화: 계약 금액, 단가, 마진율 등 기밀 제거 여부",
            "• 3. 자격증명 제거: 소스코드 내 API Key, ID/PW 제거 여부",
            "• 4. 보안 가이드 준수: 보안실 배포 규정 및 자가점검 여부",
            "• 5. 유출 위반 차단: 타 부서 및 망 분리 활용 시 유출 위험 제거 여부",
            "\n⚠️ Red Flag 경고 연동:",
            "• 하나라도 미준수 시 '보안 주의 알림'이 대시보드에 연동되며, 중대한 자격증명 노출은 즉시 불합격 처리 대상이 됩니다.",
            "\n오른쪽 화면: 실제 채점 페이지 좌측 하단 보안 체크리스트 영역"
        ],
        "media__1784176902473.png",
        "심사위원",
        "12"
    )

    # =========================================================================
    # SLIDE 13: Rubrics Weights (심사위원 - Level에 따라 채점 축과 가중치가 달라집니다)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13)
    add_two_column_slide(
        slide13,
        "Level에 따라 채점 축과 가중치가 달라집니다",
        "심사 자격 레벨별 가중치 및 배점 척도",
        [
            "자격 등급별 평가 주안점 및 비율:",
            "• Level 3 (AI Champion): 현업 적용성 및 기여 효과 검증",
            "   - 문제 해결 난이도 (30%)",
            "   - 구현 완성도 (40%)",
            "   - 확산 및 자산화 (30%)",
            "• Level 4 (AI Specialist): 심화 기술력 및 보안 통제 깊이 검증",
            "   - 기술적 구현력 (40%)",
            "   - 데이터 활용 및 보안성 (30%)",
            "   - 비즈니스 기여도 (30%)",
            "\n오른쪽 화면: 실제 채점 페이지 우측의 루브릭 점수 선택 라디오 영역"
        ],
        "media__1784176902473.png",
        "심사위원",
        "13"
    )

    # =========================================================================
    # SLIDE 14: Scoring Rules & Question Bank (심사위원 - 점수와 의견은 반드시 같은 근거를 가리켜야 합니다)
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_bg(slide14)
    add_two_column_slide(
        slide14,
        "점수와 의견은 반드시 같은 근거를 가리켜야 합니다",
        "세부 루브릭 채점 및 질문은행 면접 툴",
        [
            "루브릭 세부 채점 가이드:",
            "• 5단계 등급 선택: (탁월 100, 양호 85, 보통 70, 미흡 50, 부족 30)",
            "• 각 항목마다 정량적/정성적 등급 부여 사유 필수 기재",
            "\n⚡ Consensus Monitor (의견 불합치 방지):",
            "• 동시 배정된 위원 간 평가 편차가 15점 이상 벌어질 경우 경고 알림 노출 ➡️ 타 위원 의견 수렴 후 조율 합의문 기재",
            "\n질문은행 도우미 (Question Bank):",
            "• 지원자 과제 유형을 선택하면 면접 질문 리스트와 모범 답변 준거 가이드를 자동으로 제시하여 객관성을 보장합니다.",
            "\n오른쪽 화면: 실제 5단계 루브릭 행 카드 및 우측 하단 질문은행 패널 렌더링 화면"
        ],
        "aica_evaluation_workspace_mockup_1784167999282.png",
        "심사위원",
        "14"
    )

    # =========================================================================
    # SLIDE 15: Save & Submit (심사위원 - 제출 전 마지막 확인이 평가 품질을 좌우합니다)
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_bg(slide15)
    add_two_column_slide(
        slide15,
        "제출 전 마지막 확인이 평가 품질을 좌우합니다",
        "평가 최종 검토 및 제출 확정",
        [
            "평가 완료 및 전송 단계:",
            "• 1. 필수값 체크: 보안 체크박스 준수, 루브릭 점수 선택, 세부 의견 누락 여부 최종 갱신",
            "• 2. 타 위원 상태 검토: 타 위원의 점수 노출 없이 평가 진행률만 확인(독립적 판단 견지)",
            "• 3. 평가 제출: [평가 완료 및 제출] 버튼을 선택해 데이터베이스 기록 완료",
            "• 4. 상태 갱신 검증: 본인 대기 큐로 복귀 후 해당 과제 상태가 '완료'인지 확인",
            "\n💡 최종 판정 계산: 3인의 패널 심사위원이 최종 제출을 모두 완료하는 즉시 파이어스토어에서 평균 점수와 판정을 계산합니다.",
            "\n오른쪽 화면: 실제 채점 입력창 하단의 저장 및 최종 제출 영역"
        ],
        "media__1784176902473.png",
        "심사위원",
        "15"
    )

    # =========================================================================
    # SLIDE 16: Troubleshooting (공통 - 문제가 생기면 상태·권한·필수값부터 확인하세요)
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_bg(slide16)
    add_two_column_slide(
        slide16,
        "문제가 생기면 상태·권한·필수값부터 확인하세요",
        "자주 묻는 질문 및 트러블슈팅 조치법",
        [
            "자주 발생하는 이슈 해결 방안:",
            "• 1. 로그인 실패: 등록된 사용자 성명과 사내 이메일 주소가 정확한지 대조해 보세요. (대소문자 및 띄어쓰기 유의)",
            "• 2. 과제 등록 불가: 필수 항목(성명, 이메일, 과제설명 등) 누락 여부 및 심사위원 3인 배정 유무를 확인하세요.",
            "• 3. 점수·판정 미표시: 배정된 3인 패널 심사위원(현업, 기술, 보안) 전원이 최종 제출을 마쳤는지 확인하세요.",
            "• 4. 자료 링크 오류: 사내 보안 인트라넷 링크의 접근 권한 허용 여부 및 URL 정상 동작 여부를 사전에 검토하세요.",
            "\n💡 완료 전 4대 필수 체크 사항: 정확한 후보자 정보, 이해상충이 없는 패널 매핑, 명확한 평가 근거 작성, 보안 결격 검증 완료",
            "\n오른쪽 화면: 실제 리액트 로그인 오류 배지 렌더링 화면"
        ],
        "media__1784176101924.png",
        "공통",
        "16"
    )

    # Save presentation
    output_filename = "AICA_System_User_Manual.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'.")

if __name__ == "__main__":
    main()
