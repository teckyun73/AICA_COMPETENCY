import sys
sys.path.append(r'C:\Users\ATECCN\AppData\Roaming\Python\Python314\site-packages')
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def build_dark_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Deep Dark Slate Concept)
    BG_COLOR = RGBColor(11, 19, 41)       # #0b1329
    CARD_BG = RGBColor(17, 28, 56)        # #111c38
    CARD_BORDER = RGBColor(30, 41, 59)    # #1e293b
    CYAN_ACCENT = RGBColor(56, 189, 248)  # #38bdf8
    BLUE_ACCENT = RGBColor(96, 165, 250)  # #60a5fa
    PURPLE_ACCENT = RGBColor(192, 132, 252) # #c084fc
    TEXT_WHITE = RGBColor(255, 255, 255)   # #ffffff
    TEXT_SLATE = RGBColor(226, 232, 240)   # #e2e8f0
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94a3b8
    YELLOW_TIP = RGBColor(250, 204, 21)    # #facc15
    RED_WARN = RGBColor(248, 113, 113)     # #f87171

    FONT_NAME = 'Malgun Gothic'

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title, category="공통", slide_num_str=""):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(13)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_ACCENT
        p_cat.font.name = FONT_NAME

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(11), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = FONT_NAME

        # Divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.color.rgb = CARD_BORDER

        # Footer Left
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10), Inches(0.4))
        p_foot = footer_box.text_frame.paragraphs[0]
        p_foot.text = "AICA 자격심사 시스템 사용자 매뉴얼 · 경영혁신실"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = TEXT_MUTED
        p_foot.font.name = FONT_NAME

        # Footer Right (Page Number)
        if slide_num_str:
            num_box = slide.shapes.add_textbox(Inches(11.833), Inches(6.9), Inches(0.9), Inches(0.4))
            p_num = num_box.text_frame.paragraphs[0]
            p_num.alignment = PP_ALIGN.RIGHT
            p_num.text = slide_num_str
            p_num.font.size = Pt(10)
            p_num.font.color.rgb = TEXT_MUTED
            p_num.font.name = FONT_NAME

    def add_card_box(slide, left, top, width, height, title, bullets):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p_t = tf.paragraphs[0]
            p_t.text = title
            p_t.font.size = Pt(16)
            p_t.font.bold = True
            p_t.font.color.rgb = BLUE_ACCENT
            p_t.font.name = FONT_NAME
            p_t.space_after = Pt(10)

        first = True if not title else False
        for b in bullets:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.text = b
            p.font.size = Pt(12)
            p.font.name = FONT_NAME
            p.space_after = Pt(6)

            if "💡" in b or "팁:" in b:
                p.font.color.rgb = YELLOW_TIP
                p.font.bold = True
            elif "⚠️" in b or "불합격" in b or "경고" in b:
                p.font.color.rgb = RED_WARN
                p.font.bold = True
            elif "⚡" in b or "추천" in b:
                p.font.color.rgb = CYAN_ACCENT
                p.font.bold = True
            elif b.startswith("•") or b.startswith("1.") or b.startswith("2.") or b.startswith("3.") or b.startswith("4."):
                p.font.color.rgb = TEXT_WHITE
            else:
                p.font.color.rgb = TEXT_SLATE

    def add_screenshot_frame(slide, left, top, width, height, image_filename):
        # Outer container frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = CARD_BG
        frame.line.color.rgb = CARD_BORDER
        frame.line.width = Pt(1)

        img_path = os.path.join(r"scratch\pdf_images", image_filename)
        if not os.path.exists(img_path):
            img_path = os.path.join(r"C:\Users\ATECCN\.gemini\antigravity\brain\251e0be8-b054-44a9-82c7-d28a8df26ff8", image_filename)

        if os.path.exists(img_path):
            margin = Inches(0.1)
            x = left + margin
            y = top + margin
            w = width - (margin * 2)
            h = height - (margin * 2)
            slide.shapes.add_picture(img_path, x, y, width=w, height=h)

    def add_two_column_slide(slide, title, left_card_title, left_bullets, right_image_name, category="공통", slide_num=""):
        add_header(slide, title, category, slide_num)
        add_card_box(slide, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.8), left_card_title, left_bullets)
        add_screenshot_frame(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8), right_image_name)

    # =========================================================================
    # SLIDE 01: Title Slide (Cover Page - Deep Dark Widescreen)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p_cat = tf1.paragraphs[0]
    p_cat.text = "AICA · AI Certification for ATEC"
    p_cat.font.size = Pt(16)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN_ACCENT
    p_cat.font.name = FONT_NAME
    p_cat.space_after = Pt(14)

    p_t1 = tf1.add_paragraph()
    p_t1.text = "자격심사 시스템"
    p_t1.font.size = Pt(38)
    p_t1.font.bold = True
    p_t1.font.color.rgb = TEXT_WHITE
    p_t1.font.name = FONT_NAME

    p_t2 = tf1.add_paragraph()
    p_t2.text = "사용자 매뉴얼"
    p_t2.font.size = Pt(38)
    p_t2.font.bold = True
    p_t2.font.color.rgb = TEXT_WHITE
    p_t2.font.name = FONT_NAME
    p_t2.space_after = Pt(16)

    p_sub = tf1.add_paragraph()
    p_sub.text = "운영간사 · 심사위원용"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.font.name = FONT_NAME

    tb_foot = s1.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(6.8), Inches(0.8))
    tf_foot = tb_foot.text_frame
    p_f1 = tf_foot.paragraphs[0]
    p_f1.text = "2026. 07"
    p_f1.font.size = Pt(11)
    p_f1.font.color.rgb = TEXT_MUTED
    p_f1.font.name = FONT_NAME

    p_f2 = tf_foot.add_paragraph()
    p_f2.text = "https://aica-competency.web.app/"
    p_f2.font.size = Pt(11)
    p_f2.font.color.rgb = CYAN_ACCENT
    p_f2.font.name = FONT_NAME

    # Right Card graphic
    rcard = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.2), Inches(4.5), Inches(5.1))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = CARD_BG
    rcard.line.color.rgb = CARD_BORDER
    rcard.line.width = Pt(1.5)

    tb_rc = s1.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(4.1), Inches(4.0))
    tf_rc = tb_rc.text_frame
    tf_rc.word_wrap = True

    p_r1 = tf_rc.paragraphs[0]
    p_r1.text = "AICA"
    p_r1.font.size = Pt(40)
    p_r1.font.bold = True
    p_r1.alignment = PP_ALIGN.CENTER
    p_r1.font.color.rgb = TEXT_WHITE
    p_r1.font.name = FONT_NAME
    p_r1.space_after = Pt(24)

    p_r2 = tf_rc.add_paragraph()
    p_r2.text = "공정한 3인 패널 심사"
    p_r2.font.size = Pt(20)
    p_r2.font.bold = True
    p_r2.alignment = PP_ALIGN.CENTER
    p_r2.font.color.rgb = TEXT_WHITE
    p_r2.font.name = FONT_NAME
    p_r2.space_after = Pt(12)

    p_r3 = tf_rc.add_paragraph()
    p_r3.text = "현업 · 기술 · 보안"
    p_r3.font.size = Pt(15)
    p_r3.alignment = PP_ALIGN.CENTER
    p_r3.font.color.rgb = TEXT_MUTED
    p_r3.font.name = FONT_NAME
    p_r3.space_after = Pt(28)

    p_r4 = tf_rc.add_paragraph()
    p_r4.text = "등록 ➔ 검토 ➔ 채점 ➔ 판정"
    p_r4.font.size = Pt(16)
    p_r4.font.bold = True
    p_r4.alignment = PP_ALIGN.CENTER
    p_r4.font.color.rgb = BLUE_ACCENT
    p_r4.font.name = FONT_NAME

    # =========================================================================
    # SLIDE 02: Workflow Overview
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_two_column_slide(
        s2,
        "두 역할이 이어져야 한 건의 심사가 완료됩니다",
        "자격 심사 전체 워크플로우",
        [
            "운영간사와 심사위원 간의 원활한 업무 연계 흐름:",
            "\n1. 과제 등록:",
            "• 간사가 지원자명, 과제설명, 증빙자료 링크 및 심사예정일 등록",
            "\n2. 3인 패널 배정:",
            "• 현업, 기술, 보안별 심사위원 자동 이해상충 검증 및 매핑",
            "\n3. 독립 심사:",
            "• 각 위원이 증빙을 다차원 검증하고 5단계 루브릭 평가 제출",
            "\n4. 종합 판정:",
            "• 3인 점수 평균 및 결격 검증 후 최종 판정서 발급 및 명패 동기화"
        ],
        "page_02.png",
        "공통",
        "02"
    )

    # =========================================================================
    # SLIDE 03: Login
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_two_column_slide(
        s3,
        "성명과 사내 이메일로 로그인합니다",
        "로그인 순서 및 계정 인증",
        [
            "1. 성명(ID)에 등록된 성명을 입력",
            "• 예: 이동훈(간사), 이재윤(현업), 오찬환(기술), 이찬영(보안)",
            "\n2. 이메일(Password)에 사내 이메일 입력",
            "• 사내 그룹웨어 이메일 주소를 입력합니다.",
            "\n3. 「시스템 로그인」 선택",
            "\n4. 상단 역할 표시가 본인 권한인지 확인",
            "\n💡 입력이 없거나 일치하지 않으면 로그인 오류가 표시됩니다. 운영간사에게 등록 정보 확인을 요청하세요."
        ],
        "page_03.png",
        "공통",
        "03"
    )

    # =========================================================================
    # SLIDE 04: Roles & Completion Criteria
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_two_column_slide(
        s4,
        "먼저 내 역할과 완료 조건을 확인하세요",
        "사용자 역할 분리 및 완료 기준",
        [
            "운영간사의 핵심 임무 및 완료 조건:",
            "• 월별 자격심사 일정 관리 및 캘린더 연동",
            "• 신규 평가과제 및 후보자 등록",
            "• 3인 패널 지정 및 이해상충 자가검증",
            "• 큐·상태·점수 진척도 실시간 모니터링",
            "• 완료 신호: 3인의 심사가 완료되고 최종 판정서 생성",
            "\n심사위원의 핵심 임무 및 완료 조건:",
            "• 제출물·시연·링크 검토 및 공통 보안 체크리스트 확인",
            "• Level 3 / Level 4 루브릭 평가 및 면접 질문 작성",
            "• 완료 신호: 내 심사 상태 배지가 「완료」로 표시됨"
        ],
        "page_04.png",
        "공통",
        "04"
    )

    # =========================================================================
    # SLIDE 05: Admin Dashboard Overview
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_two_column_slide(
        s5,
        "대시보드에서 전체 심사 현황부터 확인합니다",
        "AICA 자격심사 관리 & 대기 큐 구조",
        [
            "1. 월별 자격심사 일정: AICA L3/L4 심사일정을 캘린더로 한눈에 모니터링",
            "2. 시뮬레이션 데이터 불러오기: 가상의 테스트 데이터셋(후보자 9명 및 3인 평가 결과) 복원",
            "3. 데이터 초기화: 실무 운영 착수를 위해 전체 평가 과제 및 결과를 깨끗하게 초기화",
            "4. 요약 지표: 전체, 완료, 진행 중, 대기 인원을 실시간 확인 (클릭 시 명단 모달 팝업)",
            "5. 관계사 평균: 완료된 심사를 기준으로 계열사별 평균 점수를 모니터링",
            "6. 전체 큐: 필터와 액션을 통한 개별 심사 통합 관리",
            "\n💡 대시보드 우측 상단 배지를 통해 테스트 시연용 가상 데이터 복원 및 실무용 초기화를 손쉽게 제어할 수 있습니다."
        ],
        "page_05.png",
        "운영간사",
        "05"
    )

    # =========================================================================
    # SLIDE 06: Stats Detail Modal
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_two_column_slide(
        s6,
        "대시보드에서 전체 심사 현황부터 확인합니다",
        "명단 보기 기능 (통계 모달 팝업)",
        [
            "💡 명단 보기 기능:",
            "각 지표(통계 카드) 클릭 시 해당 대상자 목록을 모달 팝업으로 즉시 확인할 수 있습니다.",
            "\n• 1. 전체 지원자 세부 정보 모달 (총 9명)",
            "• 2. 평가 완료 인원 세부 정보 모달 (2명)",
            "• 3. 평가 진행 중 인원 세부 정보 모달 (3명)",
            "• 4. 심사 대기 인원 세부 정보 모달 (4명)",
            "\n모달 내에서 관계사, 부서, 인증레벨, 제출과제명, 배정패널 및 상태를 즉시 조망할 수 있습니다."
        ],
        "page_06.png",
        "운영간사",
        "06"
    )

    # =========================================================================
    # SLIDE 07: Candidate Registration Form
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_two_column_slide(
        s7,
        "신규 과제는 필수 정보와 증빙을 한 번에 등록합니다",
        "후보자 정보 및 증빙 링크 등록 폼",
        [
            "01 지원자 정보:",
            "• 성명, 사내 이메일, 소속 관계사/부서, 자격 레벨(L3/L4), 과제 유형, 심사 예정일자",
            "\n02 과제 내용 및 6S 요약:",
            "• 평가 대상 과제명, Pain Point, AI 해결 방안, 핵심 기대효과 기술",
            "\n03 증빙 파일 링크 등록:",
            "• 디렉토리 구조 요약, 결과 보고서 PDF 주소, Git 소스코드 저장소, 시연 동영상 URL"
        ],
        "page_07.png",
        "운영간사",
        "07"
    )

    # =========================================================================
    # SLIDE 08: Panel Assignment & Conflict of Interest
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_two_column_slide(
        s8,
        "패널은 전문분야별 1인씩, 이해상충 없이 배정합니다",
        "심사위원회 패널 구성 및 최적 추천",
        [
            "1. 현업·사업성 위원 (가치성 검증): 업무 프로세스 개선율 및 사업 효과성 채점",
            "2. AI·기술 위원 (구현성 검증): 기술적 난이도 및 구현 완성도 채점",
            "3. 보안·거버넌스 위원 (신뢰성 검증): 데이터 활용 보안 및 내부 가이드 채점",
            "\n⚡ 최적 패널 자동 추천 로직:",
            "• 피평가자와 동일한 사명 소속 위원을 자동 배제하여 후보군을 우선 순위 정렬합니다.",
            "\n⚠️ 간사가 위원을 직접 수동 선택하는 경우, 소속이 일치하는 위원 옆에는 빨간색 [이해상충 주의] 레이블이 노출됩니다."
        ],
        "page_08.png",
        "운영간사",
        "08"
    )

    # =========================================================================
    # SLIDE 09: Queue Filtering & Status Monitoring
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_two_column_slide(
        s9,
        "필터와 상태값으로 병목을 빠르게 찾습니다",
        "진행 상태 분기 및 필터 모니터링",
        [
            "통합 대기 큐 및 필터 시스템 활용법:",
            "• 필터 종류: 관계사별, 인증 레벨별, 평가 진행 상태별 필터 제공",
            "• 1. 평가대기: 패널 배정 완료 후 아직 위원이 채점을 시작하지 않음",
            "• 2. 심사중: 1명 이상의 위원이 부분 채점을 기록하거나 제출함",
            "• 3. 완료: 3인 패널이 모두 최종 평가를 제출하여 종합 평균이 집계됨",
            "• 4. 보고서 보기: [결과 보기] 액션 버튼을 눌러 종합판정서 조회 및 인쇄",
            "\n💡 팁: '타 위원 진척도' 항목에서 실시간 배지 도트를 확인해 미제출 위원을 즉시 판별할 수 있습니다."
        ],
        "page_09.png",
        "운영간사",
        "09"
    )

    # =========================================================================
    # SLIDE 10: Final Verdict Report & Signatures
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_two_column_slide(
        s10,
        "3인 제출 후 종합판정서가 완성됩니다",
        "종합판정 보고서 발급 및 명패 동기화",
        [
            "종합판정서 구성 항목 및 검증 가이드:",
            "• 지원자 기본 정보 및 심사 시행/제출/승인 일자 명시",
            "• 종합 평균 점수 시각화 및 최종 합격/불합격 판정",
            "• 3인 배정 위원의 세부 루브릭 채점 점수 및 종합 정성 평가 의견 노출",
            "• 하단 서명란: 배정된 실제 심사위원 성명으로 명패 자동 변경 완료",
            "• 출력 최적화: [결과 보고서 인쇄(PDF)] 버튼을 클릭하여 A4 최적화 보고서 출력",
            "\n⚠️ 결격 사유 발생 시: 표절, 중요 기밀 및 API Key 노출 등이 발견되어 결격 처리된 경우, 평균 점수가 70점 이상이어도 최종 판정은 '불합격' 처리됩니다."
        ],
        "page_10.png",
        "운영간사",
        "10"
    )

    # =========================================================================
    # SLIDE 11: Reviewer Queue & Priority
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11)
    add_two_column_slide(
        s11,
        "내 배정 과제 큐에서 심사 우선순위를 잡습니다",
        "심사위원 배정 큐 및 월별 일정 모니터링",
        [
            "심사위원 전용 배정 과제 리스트 가이드:",
            "• 성명 및 이메일 로그인 시, 본인이 패널로 등록된 평가 과제 및 월별 일정 캘린더가 렌더링됩니다.",
            "• 테이블 컬럼: 지원자 정보, 소속 법인/부서, 자격 인증 레벨, 내 평가 진행 현황(대기/완료), 타 위원 제출 진척 상태 배지 표시",
            "• [심사하기/수정]: 채점을 시작하거나 이미 제출한 내용을 재조정할 때 버튼을 선택해 작업장(Workspace)으로 진입합니다."
        ],
        "page_11.png",
        "심사위원",
        "11"
    )

    # =========================================================================
    # SLIDE 12: Evidence Verification - Tab 1 (Report Summary)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12)
    add_two_column_slide(
        s12,
        "채점 전에 제출물의 연결성과 재현 가능성을 봅니다",
        "3대 증빙자료 다차원 교차 검증 (탭 1: 보고서 요약)",
        [
            "제출물 상세 탭(Tab) 활용 방식:",
            "• 보고서 요약 (Tab 1): 과제 제안 배경 및 Pain Point 극복 방안, 6S 요약 검토 (문제-해결-기대효과의 논리적 일관성 체크)",
            "• 제출일 및 심사 시행 일자 메타 바를 통한 일정 확인",
            "• 소스코드 구조 (Tab 2): 디렉토리 트리 구조, 라이브러리 구성 및 API Key 노출 방지 대책 검토",
            "• 시연·링크 (Tab 3): 결과 보고서 PDF 원본, Git 저장소 주소, 구동 동영상 링크 유효성 검증",
            "\n💡 소스코드 및 구동 영상 링크가 손상되었을 경우 정성 의견란에 구체적인 비작동 사유를 기재해 주세요."
        ],
        "page_12.png",
        "심사위원",
        "12"
    )

    # =========================================================================
    # SLIDE 13: Evidence Verification - Tab 2 (Code Structure)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13)
    add_two_column_slide(
        s13,
        "채점 전에 제출물의 연결성과 재현 가능성을 봅니다",
        "3대 증빙자료 다차원 교차 검증 (탭 2: 소스코드 구조)",
        [
            "소스코드 아키텍처 및 보안 검토:",
            "• 제출 소스코드 디렉토리 트리를 시각적으로 확인",
            "• `requirements.txt` 또는 `package-lock.json` 등 패키지 잠금 내역 및 `env.example` 구성도 검토",
            "• 하드코딩된 API Key, DB 비밀번호, 접근 토큰 포함 여부 집중 조명",
            "• PII 정규식 필터 및 예외 처리 로직 구현 확인"
        ],
        "page_13.png",
        "심사위원",
        "13"
    )

    # =========================================================================
    # SLIDE 14: Evidence Verification - Tab 3 (Demo & Links)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14)
    add_two_column_slide(
        s14,
        "채점 전에 제출물의 연결성과 재현 가능성을 봅니다",
        "3대 증빙자료 다차원 교차 검증 (탭 3: 시연 & 링크)",
        [
            "16:9 데모 샌드박스 플레이어 & 라이브 미리보기:",
            "• 16:9 표준 비율 창: 뷰포트 규격에 최적화된 16대9 디스플레이 플레이어 탑재",
            "• 라이브 데모 연동 미리보기: 심사 창 내부 뷰포트(iframe)에서 데모 웹 앱을 직접 가동 및 미리보기",
            "• MP4 동영상 재생 모드: 시연 MP4 동영상 플레이어 및 타임라인 슬라이더 연동",
            "\n과제 증빙용 하이퍼링크 3종 검증:",
            "• 1. 결과 보고서 원본 (그룹웨어 PDF)",
            "• 2. 코드 저장소 (GitLab Repository)",
            "• 3. 실제 구동 데모 (사내 Sandbox URL)"
        ],
        "page_14.png",
        "심사위원",
        "14"
    )

    # =========================================================================
    # SLIDE 15: Security Checklist & Red Flag
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_bg(s15)
    add_two_column_slide(
        s15,
        "공통 보안 체크리스트는 모든 심사의 선행 조건입니다",
        "5대 공통 보안성 자가진단 규칙",
        [
            "모든 심사의 최우선 선행 체크리스트:",
            "• 1. 개인정보(PII): 이름, 메일주소, 연락처 등 완벽 마스킹 여부",
            "• 2. 기밀 정보 익명화: 계약 금액, 단가, 마진율 등 기밀 제거 여부",
            "• 3. 자격증명 제거: 소스코드 내 API Key, ID/PW 제거 여부",
            "• 4. 보안 가이드 준수: 보안실 배포 규정 및 자가점검 여부",
            "• 5. 유출 위반 차단: 타 부서 및 망 분리 활용 시 유출 위험 제거 여부",
            "\n⚠️ Red Flag 경고 연동: 하나라도 미준수 시 '보안 주의 알림'이 대시보드에 연동되며, 중대한 자격증명 노출은 즉시 불합격 처리 대상이 됩니다."
        ],
        "page_15.png",
        "심사위원",
        "15"
    )

    # =========================================================================
    # SLIDE 16: Rubric Weights by Level
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_bg(s16)
    add_two_column_slide(
        s16,
        "Level에 따라 채점 축과 가중치가 달라집니다",
        "심사 자격 레벨별 가중치 및 배점 척도",
        [
            "자격 등급별 평가 주안점 및 비율:",
            "• Level 3 (AI Champion): 현업 적용성 및 기여 효과 검증",
            "   - 문제 해결 난이도 (30%)",
            "   - 구현 완성도 (40%)",
            "   - 확산 및 자산화 (30%)",
            "• Level 4 (AI Specialist): 심화 기술력 및 보안 통제 깊이 검증",
            "   - 기술적 구현력 (40%)",
            "   - 데이터 활용 및 보안성 (30%)",
            "   - 비즈니스 기여도 (30%)"
        ],
        "page_16.png",
        "심사위원",
        "16"
    )

    # =========================================================================
    # SLIDE 17: Rubric Scoring & Question Bank
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    set_bg(s17)
    add_two_column_slide(
        s17,
        "점수와 의견은 반드시 같은 근거를 가리켜야 합니다",
        "세부 루브릭 채점 및 질문은행 면접 툴",
        [
            "루브릭 세부 채점 가이드:",
            "• 5단계 등급 선택: (탁월 100, 양호 85, 보통 70, 미흡 50, 부족 30)",
            "• 각 항목마다 정량적/정성적 등급 부여 사유 필수 기재",
            "\n⚡ Consensus Monitor (의견 불합치 방지):",
            "• 동시 배정된 위원 간 평가 편차가 15점 이상 벌어질 경우 경고 알림 노출 ➡️ 타 위원 의견 수렴 후 조율 합의문 기재",
            "\n질문은행 도우미 (Question Bank):",
            "• 지원자 과제 유형을 선택하면 면접 질문 리스트와 모범 답변 준거 가이드를 자동으로 제시하여 객관성을 보장합니다."
        ],
        "page_17.png",
        "심사위원",
        "17"
    )

    # =========================================================================
    # SLIDE 18: Final Review & Submission
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    set_bg(s18)
    add_two_column_slide(
        s18,
        "제출 전 마지막 확인이 평가 품질을 좌우합니다",
        "평가 최종 검토 및 제출 확정",
        [
            "평가 완료 및 전송 단계:",
            "• 1. 필수값 체크: 보안 체크박스 준수, 루브릭 점수 선택, 세부 의견 누락 여부 최종 갱신",
            "• 2. 타 위원 상태 검토: 타 위원의 점수 노출 없이 평가 진행률만 확인(독립적 판단 견지)",
            "• 3. 평가 제출: [평가 완료 및 제출] 버튼을 선택해 데이터베이스 기록 완료",
            "• 4. 상태 갱신 검증: 본인 대기 큐로 복귀 후 해당 과제 상태가 '완료'인지 확인",
            "\n💡 최종 판정 계산: 3인의 패널 심사위원이 최종 제출을 모두 완료하는 즉시 파이어스토어에서 평균 점수와 판정을 계산합니다."
        ],
        "page_18.png",
        "심사위원",
        "18"
    )

    # =========================================================================
    # SLIDE 19: Troubleshooting & FAQ
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    set_bg(s19)
    add_two_column_slide(
        s19,
        "문제가 생기면 상태·권한·필수값부터 확인하세요",
        "자주 묻는 질문 및 트러블슈팅 조치법",
        [
            "자주 발생하는 이슈 해결 방안:",
            "• 1. 로그인 실패: 등록된 사용자 성명과 사내 이메일 주소가 정확한지 대조해 보세요.",
            "• 2. 과제 등록 불가: 필수 항목(성명, 이메일, 과제설명 등) 누락 여부 및 심사위원 3인 배정 유무를 확인하세요.",
            "• 3. 점수·판정 미표시: 배정된 3인 패널 심사위원(현업, 기술, 보안) 전원이 최종 제출을 마쳤는지 확인하세요.",
            "• 4. 자료 링크 오류: 사내 보안 인트라넷 링크의 접근 권한 허용 여부 및 URL 정상 동작 여부를 사전에 검토하세요.",
            "\n💡 완료 전 4대 필수 체크 사항: 정확한 후보자 정보, 이해상충이 없는 패널 매핑, 명확한 평가 근거 작성, 보안 결격 검증 완료"
        ],
        "page_19.png",
        "공통",
        "19"
    )

    # =========================================================================
    # SLIDE 20: End of Document
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    set_bg(s20)

    tb20 = s20.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
    tf20 = tb20.text_frame
    tf20.word_wrap = True

    p_e1 = tf20.paragraphs[0]
    p_e1.alignment = PP_ALIGN.CENTER
    p_e1.text = "End of Document"
    p_e1.font.size = Pt(36)
    p_e1.font.bold = True
    p_e1.font.color.rgb = TEXT_WHITE
    p_e1.font.name = FONT_NAME
    p_e1.space_after = Pt(16)

    p_e2 = tf20.add_paragraph()
    p_e2.alignment = PP_ALIGN.CENTER
    p_e2.text = "AICA 자격심사 시스템 사용자 매뉴얼 · 경영혁신실"
    p_e2.font.size = Pt(16)
    p_e2.font.color.rgb = CYAN_ACCENT
    p_e2.font.name = FONT_NAME

    # Save Presentation
    out_pptx = "AICA_System_User_Manual.pptx"
    prs.save(out_pptx)
    print(f"Presentation saved successfully as '{out_pptx}'.")

if __name__ == "__main__":
    build_dark_presentation()
